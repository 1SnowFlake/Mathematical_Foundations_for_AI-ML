from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define custom colors
    brand_primary = RGBColor(99, 102, 241) # Indigo
    brand_text = RGBColor(28, 25, 23)      # Dark Gray
    brand_subtext = RGBColor(120, 113, 108)
    
    # Helper to add a title slide
    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = brand_primary
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = brand_text
        return slide

    # Helper to add a content slide
    def add_content_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = brand_primary
        
        tf = body.text_frame
        tf.clear() # Clear default
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = brand_text
            p.space_after = Pt(10)
        return slide

    # Slide 1: Title
    add_title_slide(
        "Interactive AI / Math / ML Syllabus",
        "A Unified Platform for Mathematical Foundations"
    )

    # Slide 2: The Problem
    add_content_slide(
        "The Problem with Traditional Learning",
        [
            "Math and ML are highly visual, but taught via static equations.",
            "Existing tools are fragmented (e.g. CNN Explainer only for CNNs, Seeing Theory only for Stats).",
            "Lack of unified, step-by-step interactive platforms bridging Theory and Practice."
        ]
    )

    # Slide 3: The Framework
    add_content_slide(
        "Our Framework & Architecture",
        [
            "Built on Next.js 16, React 19, TypeScript, and Tailwind v4.",
            "Layered Component Architecture:",
            "  1. Routing (App Router, Dynamic Sidebar)",
            "  2. Primitives (VectorCanvas, GraphEditor, Surface3D, MathBlock, ParamPanel)",
            "  3. Algorithms (graphAlgorithms.ts, stats.ts, matrix.ts)",
            "Everything is 'Textbook Meets Playground' - theory and interactive widgets coexist."
        ]
    )

    # Slide 4: How It Works
    add_content_slide(
        "How It Works",
        [
            "Select from 37 topics across 10 sections (Linear Algebra -> Transformers).",
            "Read theory with LaTeX-rendered math (KaTeX).",
            "Interact directly with widgets (drag vectors, edit graphs, step through algorithms).",
            "VCR-style controls (Play, Pause, Step) for algorithms like BFS, DFS, Minimax.",
            "Solve hands-on 'Try It Yourself' challenges on each page."
        ]
    )

    # Slide 5: UI/UX Design
    add_content_slide(
        "UI Design Principles",
        [
            "Dark Mode & Responsive: Seamless theme switching, touch-friendly.",
            "Modern Aesthetics: Animated gradients, smooth state transitions, glassmorphism elements.",
            "Progress Dashboard: Circular progress rings, topic checkboxes, local tracking.",
            "Accessible: Semantic HTML, high contrast text, visual cues beyond color."
        ]
    )

    # Slide 6: Competitive Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    title = slide.shapes.title
    title.text = "Competitive Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = brand_primary
    
    rows, cols = 6, 4
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9.0)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Feature", "Our Project", "TF Playground", "Seeing Theory"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = brand_primary
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        
    data = [
        ["Scope", "Full ML/Math Syllabus", "Neural Nets Only", "Probability Only"],
        ["Interactivity", "Drag, Step-Through", "Drag + Train", "Click-based"],
        ["Algorithm Playback", "BFS, DFS, Minimax...", "1 Model", "None"],
        ["Progress Tracking", "Yes (Dashboard)", "No", "No"],
        ["Mobile Friendly", "Yes", "Desktop-First", "No"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_data
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = brand_text

    # Slide 7: Interactive ML Algorithms
    add_content_slide(
        "Showcasing ML Visualizations",
        [
            "Vectors & Operations (Addition, Dot/Cross Product).",
            "Graph Theory (BFS, DFS, step-by-step execution).",
            "Neural Networks (Perceptron decision boundaries).",
            "Optimization (3D Surface Gradient Descent paths).",
            "Probability (Dynamic interactive distributions).",
            "Game Playing (Minimax & Alpha-Beta Pruning)."
        ]
    )

    # Slide 8: Future Roadmap
    add_content_slide(
        "Future Enhancements",
        [
            "Backend user authentication and cloud progress syncing.",
            "Interactive quizzes and graded assessments.",
            "More advanced architectures (GANs, LSTMs in real-time).",
            "Exportable interactive notebooks for educators."
        ]
    )

    prs.save('Project_Presentation.pptx')
    print("Presentation created successfully as 'Project_Presentation.pptx'")

if __name__ == '__main__':
    create_presentation()
